from django.shortcuts import render
from .forms import ImageForm
from .models import Image,PDFFile,imagefile
from django.shortcuts import render, redirect,get_object_or_404
from django.http import FileResponse
from docx2pdf import convert
from .forms import Wordform,CompressedPDFForm
import os
from django.core.files.storage import FileSystemStorage
from docx import Document
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import PyPDF2
from django.http import HttpResponse
from wsgiref.util import FileWrapper
from PyPDF2 import PdfReader, PdfWriter,PdfFileMerger,PdfFileReader, PdfFileWriter
from .forms import PDFFileForm
from django.utils.text import slugify
from django.core.files import File
from pdf2image import convert_from_path
from .forms import imageFileForm
from .converter import convert_ppt_to_pdf
import zipfile
from django.conf import settings
from pptx import Presentation
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, PageBreak, Paragraph
from reportlab.platypus.tables import Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from fpdf import FPDF
def home(request):
 if request.method == "POST":
  form = ImageForm(request.POST, request.FILES)
  if form.is_valid():
   form.save()
 form = ImageForm()
 img = Image.objects.all()
 return render(request, 'myapp/home.html', {'img':img, 'form':form})

def convert_word_to_pdf(request):
    if request.method == 'POST':
        form = Wordform(request.POST, request.FILES)
        if form.is_valid():
            uploaded_file = form.save(commit=False)
            uploaded_file.user = request.user
            uploaded_file.save()

            # Read the uploaded Word document
            document = Document(uploaded_file.file.path)

            # Create a BytesIO object to hold the PDF content
            pdf_buffer = BytesIO()

            # Create a PDF from the Word document using reportlab
            pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
            for paragraph in document.paragraphs:
                pdf.drawString(100, 700, paragraph.text)
            pdf.save()

            # Reset the buffer position
            pdf_buffer.seek(0)

            # Provide the PDF file for download with a custom filename
            response = FileResponse(pdf_buffer, as_attachment=True, filename="converted_document.pdf")
            return response

    else:
        form = Wordform()
    return render(request, 'myapp/word.html', {'form': form})

def merge_pdfs(request):
    
    if request.method == 'POST':
        form = PDFFileForm(request.POST, request.FILES)
        if form.is_valid():
            # Save the uploaded PDF file
            pdf_file = form.save()

            # Merge PDF files
            merger = PdfFileMerger()
            for pdf in PDFFile.objects.all():
                merger.append(pdf.pdf_file.path)
            merged_pdf = merger.write("merged_all_pages.pdf")
            merger.close()

            return redirect('pdf_merge:merged_pdf')
    else:
        form = PDFFileForm()

    return render(request, 'myapp/merge_pdfs.html', {'form': form})
def compress_pdf(request):
    if request.method == 'POST':
        form = CompressedPDFForm(request.POST, request.FILES)
        if form.is_valid():
           
            original_pdf = form.save()


            compressed_pdf_writer = PdfFileWriter()

            
            with open(original_pdf.original_pdf.path, 'rb') as pdf_file:
                original_pdf_reader = PdfFileReader(pdf_file)

                
                for page_num in range(original_pdf_reader.getNumPages()):
                    page = original_pdf_reader.getPage(page_num)
                    page.compressContentStreams()  
                    compressed_pdf_writer.addPage(page)

                
                compressed_pdf_filename = f"compressed_{slugify(original_pdf.original_pdf.name.split('/')[-1])}.pdf"
                compressed_pdf_path = os.path.join(settings.MEDIA_ROOT, 'compressed_pdfs', compressed_pdf_filename)
                with open(compressed_pdf_path, 'wb') as compressed_pdf_file:
                    compressed_pdf_writer.write(compressed_pdf_file)

            # Save the path of the compressed PDF
            original_pdf.compressed_pdf = 'compressed_pdfs/' + compressed_pdf_filename
            original_pdf.save()

            # Serve the compressed PDF as a downloadable attachment
            compressed_pdf_file = open(compressed_pdf_path, 'rb')
            response = FileResponse(compressed_pdf_file)
            response['Content-Disposition'] = f'attachment; filename="{compressed_pdf_filename}"'
            return response
    else:
        form = CompressedPDFForm()

    return render(request, 'myapp/compress_pdf.html', {'form': form})

def pdf_to_images(request):
    if request.method == 'POST':
        form =  imageFileForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_file = form.save()

            # Convert the uploaded PDF to images
            images = convert_from_path(pdf_file.pdf.path)

            # Use the PDF file name as the title
            pdf_file_name = os.path.basename(pdf_file.pdf.name)
            title = os.path.splitext(pdf_file_name)[0]

            # Create a directory to store the images
            image_dir = os.path.join(settings.MEDIA_ROOT, 'images', title)
            os.makedirs(image_dir, exist_ok=True)

            # Save the images to the directory
            image_paths = []
            for i, image in enumerate(images):
                image_name = f'image_{i}.jpg'
                image_path = os.path.join(image_dir, image_name)
                image.save(image_path, 'JPEG')
                image_paths.append(image_path)

            # Create a ZIP archive of the images
            zip_file_path = os.path.join(settings.MEDIA_ROOT, 'images', f'{title}.zip')
            with zipfile.ZipFile(zip_file_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for image_path in image_paths:
                    zipf.write(image_path, os.path.relpath(image_path, image_dir))

            # Provide a download link for the ZIP archive
            response = HttpResponse(open(zip_file_path, 'rb').read(), content_type='application/zip')
            response['Content-Disposition'] = f'attachment; filename="{title}.zip"'
            return response

    else:
        form =  imageFileForm()

    return render(request, 'myapp/pdf_to_images.html', {'form': form})

from django.http import HttpResponse
from django.shortcuts import render
from django import forms
from pptx import Presentation
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os

class PPTToPDFForm(forms.Form):
    ppt_file = forms.FileField()

def convert_ppt(request):
    pdf_url = None

    if request.method == 'POST':
        form = PPTToPDFForm(request.POST, request.FILES)

        if form.is_valid():
            ppt_file = form.cleaned_data['ppt_file']
            ppt = Presentation(ppt_file)

            pdf_buffer = BytesIO()
            pdf = canvas.Canvas(pdf_buffer, pagesize=letter)

            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        pdf.drawString(100, 500, text)

                    if hasattr(shape, 'image'):
                        image = shape.image
                        image_data = image.blob
                        image_stream = BytesIO(image_data)
                        pdf.drawImage(image_stream, 100, 100, width=400, height=400)

                pdf.showPage()

            pdf.save()
            pdf_buffer.seek(0)

            response = HttpResponse(pdf_buffer, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{os.path.splitext(ppt_file.name)[0]}.pdf"'
            return response
    else:
        form = PPTToPDFForm()

    return render(request, 'myapp/convert_ppt.html', {'form': form, 'pdf_url': pdf_url})

from django.shortcuts import render, HttpResponse
import img2pdf
from PIL import Image


def imgtopdf(request):
    if request.method == 'POST':
        img = request.FILES['img']
        pdf = img2pdf.convert(img)

        return HttpResponse(pdf, content_type='application/pdf')
    return render(request, 'myapp/image_to_pdf.html')
def jpgtopng(request):
    if request.method == 'POST':
        img = request.FILES['img']
        im = Image.open(img)
        # im.save('Foto.png')
        response = HttpResponse(content_type="image/png")
        im.save(response, "PNG")
        return response
    return render(request, 'myapp/jpgtopng.html')


def png(request):
    if request.method == 'POST':
        convertedimage = request.FILES['img']
        type = request.POST.get('convert-type')
        if type == 'jpg':
            im = Image.open(convertedimage)
            response = HttpResponse(content_type="image/JPEG")
            im.convert('RGB').save(response, "JPEG")
            return response
        elif type == 'webp':
            im = Image.open(convertedimage)
            response = HttpResponse(content_type="image/png")
            im.convert('RGB').save(response, "webp")
            return response

        return HttpResponse('not jpg')
    return render(request, 'myapp/png.html')

from django.shortcuts import render, redirect
from django.http import FileResponse
from .models import UploadedFile
from .forms import UploadFileForm
import os
import tempfile
from docx import Document
import PyPDF2

from django.conf import settings  # Import the Django settings module

# ...

from django.conf import settings
from django.conf import settings
from urllib.parse import urljoin  # Use urljoin to construct the URL
from django.conf import settings
from urllib.parse import urljoin
def handle_file(request):
    if request.method == 'POST':
        form = UploadFileForm(request.POST, request.FILES)
        if form.is_valid():
            instance = form.save()

            with tempfile.TemporaryDirectory() as tmpdir:
                pdf_path = instance.pdf_file.path
                pdf_reader = PyPDF2.PdfFileReader(open(pdf_path, 'rb'))
                doc = Document()
                for page_num in range(pdf_reader.numPages):
                    page = pdf_reader.getPage(page_num)
                    text = page.extractText()
                    doc.add_paragraph(text)

                # Set the path to save the Word file within the MEDIA_ROOT directory
                word_path = os.path.join(settings.MEDIA_ROOT, 'converted.docx')
                doc.save(word_path)

                # Check if the Word file was successfully created
                if os.path.exists(word_path):
                    # Set the FieldFile to the saved Word file path
                    instance.word_file = 'converted.docx'
                    instance.save()

                    # Calculate the full URL for the Word file using the FieldFile's URL attribute
                    word_url = urljoin(settings.MEDIA_URL, instance.word_file.url)

                    # Return a response with the download link
                    return render(request, 'myapp/download.html', {'word_url': word_url})
                else:
                    # Handle the case where the conversion failed
                    # You can return an error message or redirect as needed
                    return render(request, 'conversion_error.html')

    else:
        form = UploadFileForm()
    return render(request, 'myapp/upload.html', {'form': form})
from django.shortcuts import render, redirect
from .forms import ConversionForm
from .models import Conversion
from django.http import HttpResponse, FileResponse
import os
import tempfile
import io
from pdf2image import convert_from_bytes
from PIL import Image
import xlsxwriter

def convert_pdf_to_excel(request):
    if request.method == 'POST':
        form = ConversionForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_file = form.cleaned_data['pdf_file']

            with tempfile.TemporaryDirectory() as temp_dir:
                pdf_bytes = pdf_file.read()
                pdf_images = convert_from_bytes(pdf_bytes)

                # Create a new directory to store image files
                images_dir = os.path.join(temp_dir, 'images')
                os.makedirs(images_dir)

                image_paths = []
                for index, image in enumerate(pdf_images, start=1):
                    image_path = os.path.join(images_dir, f'page_{index}.png')
                    image.save(image_path)
                    image_paths.append(image_path)

                # Create a new Excel file
                excel_path = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False).name

                # Create an XlsxWriter workbook and add a worksheet.
                workbook = xlsxwriter.Workbook(excel_path)
                worksheet = workbook.add_worksheet()

                for index, image_path in enumerate(image_paths, start=1):
                    # Add each image to the Excel sheet
                    worksheet.insert_image(f'A{index}', image_path, {'x_offset': 10, 'y_offset': 10})

                workbook.close()

                # Create a new Conversion record in the database
                conversion = Conversion()
                conversion.pdf_file = pdf_file
                with open(excel_path, 'rb') as excel_file:
                    conversion.excel_file.save('converted_excel.xlsx', File(excel_file))

                conversion.save()

                response = FileResponse(open(excel_path, 'rb'))
                response['Content-Type'] = 'application/ms-excel'
                response['Content-Disposition'] = f'attachment; filename="converted_excel.xlsx"'
                return response

    else:
        form = ConversionForm()

    return render(request, 'myapp/uploadexel.html', {'form': form})

# views.py
from django.shortcuts import render, redirect
from .forms import PDFFileForm
from PyPDF2 import PdfFileWriter, PdfFileReader
import os

def set_password_to_pdf(request):
    if request.method == 'POST':
        form = PDFFileForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_file = form.save(commit=False)
            pdf_file.save()

            input_pdf = pdf_file.pdf_file.path
            output_pdf = os.path.join('media/pdf_files/', os.path.basename(input_pdf))

            pdf = PdfFileReader(open(input_pdf, 'rb'))
            pdf_writer = PdfFileWriter()

            for page_num in range(pdf.numPages):
                page = pdf.getPage(page_num)
                pdf_writer.addPage(page)

            pdf_writer.encrypt(pdf_file.password)

            with open(output_pdf, 'wb') as output_file:
                pdf_writer.write(output_file)

            pdf_file.pdf_file.name = output_pdf
            pdf_file.save()

            return redirect('pdf_list')

    else:
        form = PDFFileForm()

    return render(request, 'myapp/set_password.html', {'form': form})



from django.shortcuts import render
from .forms import CompressedPDFForm

def your_view(request):
    if request.method == 'POST':
        form = CompressedPDFForm(request.POST, request.FILES)
        if form.is_valid():
            # Access the uploaded file
            uploaded_file = request.FILES['original_pdf']

            # Process or save the file
            # For example, you can save it to a specific location:
            with open('path/to/your/uploads/' + uploaded_file.name, 'wb') as destination:
                for chunk in uploaded_file.chunks():
                    destination.write(chunk)

            # Continue with other processing or redirect as needed

    else:
        form = CompressedPDFForm()

    return render(request, 'compress_pdf.html', {'form': form})




