# myapp/converter.py
import os
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from django.conf import settings

def convert_ppt_to_pdf(ppt_file_path):
    pdf_buffer = os.path.join(settings.MEDIA_ROOT, 'pdfs', 'result.pdf')
    ppt = Presentation(ppt_file_path)
    pdf = canvas.Canvas(pdf_buffer, pagesize=letter)

    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.drawString(100, 600, shape.text)

        pdf.showPage()

    pdf.save()
    return pdf_buffer
